import os
import fnmatch
import time
import csv
from functools import wraps
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import pythoncom
from win32com.client import Dispatch
import xlrd
from pymongo import MongoClient
from dotenv import load_dotenv
from nltk.tokenize import word_tokenize
from nltk.corpus import stopwords
from nltk.stem import WordNetLemmatizer
from PyPDF2 import PdfReader
from wordsegment import load as load_wordsegment, segment
from symspellpy import SymSpell, Verbosity
from nltk.corpus import words as nltk_words
import re

def timer(func):
    @wraps(func)
    def wrapper(*args, **kwargs):
        start_time = time.time()
        result = func(*args, **kwargs)
        end_time = time.time()
        duration = end_time - start_time
        print(f"\nTime taken by {func.__name__}: {duration:.2f} seconds")
        return result
    return wrapper

def search_files_os_walk(root_dir, file_types):
    matches = []
    for root, dirnames, filenames in os.walk(root_dir):
        for file_type in file_types:
            for filename in fnmatch.filter(filenames, file_type):
                matches.append(os.path.join(root, filename))
    return matches

def read_doc_file(file_path):
    pythoncom.CoInitialize()
    word = Dispatch("Word.Application")
    doc = word.Documents.Open(file_path)
    file_data = ' '.join([para.Range.Text for para in doc.Paragraphs])
    doc.Close()
    word.Quit()
    return file_data

def clean_pdf_text(file_path: str) -> str:
    """
    Reads a PDF file, removes stop words, lemmatizes text, and returns cleaned text.

    Args:
        file_path (str): Path to the PDF file.

    Returns:
        str: Cleaned text from the PDF.
    """
    # Initialize lemmatizer, SymSpell, and WordSegment
    lemmatizer = WordNetLemmatizer()
    sym_spell = SymSpell(max_dictionary_edit_distance=2, prefix_length=7)
    sym_spell.load_dictionary("frequency_dictionary_en_82_765.txt", term_index=0, count_index=1)
    load_wordsegment()  # Load WordSegment model

    # Load stop words and English words
    stop_words = set(stopwords.words('english'))
    english_words = set(nltk_words.words())

    try:
        with open(file_path, 'rb') as pdf_file:
            reader = PdfReader(pdf_file)
            text = ''

            # Extract text from all pages
            for page in reader.pages:
                page_text = page.extract_text() + ' '
                text += page_text

            # Clean and tokenize text
            text = re.sub(r'[^a-zA-Z\s]', '', text)
            tokens = word_tokenize(text)

            corrected_tokens = []
            for token in tokens:
                # Skip stop words
                if token.lower() in stop_words:
                    continue

                # Use WordSegment to split concatenated words
                segmented_words = segment(token)

                for word in segmented_words:
                    # Spell correction with SymSpell
                    suggestions = sym_spell.lookup(word, Verbosity.CLOSEST, max_edit_distance=2)
                    corrected_word = suggestions[0].term if suggestions else word
                    # Lemmatize the corrected word
                    lemmatized_word = lemmatizer.lemmatize(corrected_word)
                    # Remove 2 or 3 letter words not in the English dictionary, except "ai"
                    if len(lemmatized_word) > 3 or lemmatized_word in english_words or lemmatized_word == "ai":
                        corrected_tokens.append(lemmatized_word)

            # Combine cleaned tokens into a single string
            cleaned_text = ' '.join(corrected_tokens)
            return cleaned_text
    except Exception as e:
        print(f"Error reading PDF {file_path}: {e}")
        return ""

def get_file_info(file_path):
    """
    Extracts file information and processes the file content based on its type.

    Args:
        file_path (str): Path to the file.

    Returns:
        tuple: File name, file type, file size, and file data.
    """
    file_name = os.path.basename(file_path)
    file_type = os.path.splitext(file_path)[1]
    file_size = os.path.getsize(file_path)
    file_data = ""

    try:
        if file_type.lower() == ".doc":
            file_data = read_doc_file(file_path)
        elif file_type.lower() == ".docx":
            doc = Document(file_path)
            file_data = ' '.join([para.text for para in doc.paragraphs])
        elif file_type.lower() == ".xls":
            wb = xlrd.open_workbook(file_path)
            file_data = ' '.join([str(cell.value) for sheet in wb.sheets() for row in range(sheet.nrows) for cell in sheet.row(row)])
        elif file_type.lower() == ".xlsx":
            wb = load_workbook(file_path, read_only=True, data_only=True)
            file_data = ' '.join([str(cell.value) for sheet in wb for row in sheet.iter_rows() for cell in row])
        elif file_type.lower() in [".ppt", ".pptx"]:
            prs = Presentation(file_path)
            file_data = ' '.join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
        elif file_type.lower() == ".pdf":
            file_data = clean_pdf_text(file_path)  # Use the new PDF cleaning function
        else:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                file_data = file.read().replace('\n', ' ').replace('\r', ' ')
    except Exception as e:
        print(f"Error reading {file_path}: {e}")

    return file_name, file_type, file_size, file_data

@timer
def write_file_info_to_csv(file_paths, csv_file_path):
    with open(csv_file_path, 'w', newline='', encoding='utf-8') as csvfile:
        csv_writer = csv.writer(csvfile)
        csv_writer.writerow(['File Name', 'File Type', 'File Size', 'File Data'])
        for file_path in file_paths:
            try:
                file_info = get_file_info(file_path)
                csv_writer.writerow(file_info)
            except Exception as e:
                print(f"Error processing {file_path}: {e}")
    print(f"\nFile information has been written to {csv_file_path}")

# Load environment variables from .env file
load_dotenv()

# Get MongoDB URI from environment variables
mongo_uri = os.getenv("MONGO_URI")
db_name = "DB-Files"
collection_name = "Files-Data"

def write_file_info_to_mongodb(file_paths, mongo_uri, db_name, collection_name):
    """
    Writes file information to a MongoDB database and tracks file formats processed and errors.
    """
    client = MongoClient(mongo_uri)
    db = client[db_name]
    collection = db[collection_name]

    processed_formats = set()
    error_formats = set()

    for file_path in file_paths:
        try:
            file_info = get_file_info(file_path)
            document = {
                "file_name": file_info[0],
                "file_type": file_info[1],
                "file_size": file_info[2],
                "file_data": file_info[3]
            }
            collection.insert_one(document)
            processed_formats.add(file_info[1].lower())  # Track successfully processed formats
        except Exception as e:
            print(f"Error processing {file_path}: {e}")
            error_formats.add(os.path.splitext(file_path)[1].lower())  # Track formats causing errors

    print("\nProcessed file formats:")
    for fmt in processed_formats:
        print(fmt)

    print("\nFile formats that caused errors:")
    for fmt in error_formats:
        print(fmt)

    print(f"\nFile information has been written to MongoDB collection '{collection_name}' in database '{db_name}'.")

@timer
def main():
    # Define the root directory to start the search
    root_directory = os.path.expanduser("~\\Downloads")
    
    # Define the file types to search for
    file_types = [
        "*.txt", "*.md",  # Existing file types
        "*.pdf",          # PDF files
        "*.doc", "*.docx",# Microsoft Word files
        "*.xls", "*.xlsx",# Microsoft Excel files
        "*.ppt", "*.pptx" # Microsoft PowerPoint files
    ]
    
    # Search for files using os.walk and fnmatch
    found_files_os_walk = search_files_os_walk(root_directory, file_types)
    
    # Print the found files using os.walk and fnmatch
    print("Files found using os.walk and fnmatch:")
    for file in found_files_os_walk:
        print(file)
    
    # Write file information to MongoDB and track formats
    write_file_info_to_mongodb(found_files_os_walk, mongo_uri, db_name, collection_name)

if __name__ == "__main__":
    main()